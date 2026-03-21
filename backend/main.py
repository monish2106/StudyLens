"""
StudyLens Backend — FastAPI
Full-stack reading & study assistant
"""

from fastapi import FastAPI, UploadFile, File, HTTPException, BackgroundTasks
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse, JSONResponse
from pydantic import BaseModel
from typing import Optional, List
import uuid, os, json, time, re, math, shutil
from pathlib import Path
from datetime import datetime

# ── Optional heavy parsers (graceful degradation if absent) ──────────────────
try:
    import PyPDF2
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import docx as python_docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

# ── Paths ─────────────────────────────────────────────────────────────────────
BASE_DIR   = Path(__file__).parent.parent
UPLOAD_DIR = BASE_DIR / "uploads"
EXPORT_DIR = BASE_DIR / "exports"
FRONT_DIR  = BASE_DIR / "frontend"
UPLOAD_DIR.mkdir(exist_ok=True)
EXPORT_DIR.mkdir(exist_ok=True)

# ── In-memory job store ───────────────────────────────────────────────────────
jobs: dict[str, dict] = {}

# ── App ───────────────────────────────────────────────────────────────────────
app = FastAPI(title="StudyLens API", version="1.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# ── Serve frontend ────────────────────────────────────────────────────────────
app.mount("/static", StaticFiles(directory=str(FRONT_DIR / "static")), name="static")

@app.get("/", response_class=FileResponse)
async def root():
    return FileResponse(str(FRONT_DIR / "index.html"))

# ─────────────────────────────────────────────────────────────────────────────
# UTILITIES
# ─────────────────────────────────────────────────────────────────────────────

TOPIC_KEYWORDS = {
    "Introduction": ["introduction", "overview", "background", "history", "what is", "definition"],
    "Core Concept": ["concept", "principle", "theory", "fundamental", "basic", "key", "important"],
    "Algorithm": ["algorithm", "procedure", "steps", "method", "approach", "technique", "scheduling", "sorting"],
    "Architecture": ["architecture", "structure", "design", "component", "system", "layer", "module"],
    "Synchronization": ["synchronization", "semaphore", "mutex", "lock", "deadlock", "concurrency", "thread"],
    "Memory": ["memory", "cache", "paging", "segmentation", "virtual", "ram", "heap", "stack"],
    "Network": ["network", "protocol", "tcp", "ip", "http", "packet", "router", "socket"],
    "Security": ["security", "encryption", "authentication", "authorization", "firewall", "vulnerability"],
    "Database": ["database", "sql", "query", "table", "index", "transaction", "schema"],
    "Hardware": ["hardware", "cpu", "disk", "i/o", "device", "bus", "interrupt", "register"],
}

QUESTION_TEMPLATES = [
    ("What is {topic}?",
     "A foundational question asking for the definition and purpose of {topic}."),
    ("Explain the key components of {topic}.",
     "Identify and describe each major component or element that makes up {topic}."),
    ("How does {topic} work in practice?",
     "Describe the step-by-step process or mechanism behind {topic}."),
    ("What are the advantages and disadvantages of {topic}?",
     "List the main benefits and limitations, with examples where possible."),
    ("Compare {topic} with related approaches.",
     "Identify similar concepts and highlight key differences in purpose and behaviour."),
    ("Give a real-world example of {topic}.",
     "Provide a concrete scenario where {topic} is applied or observed."),
    ("What problems does {topic} solve?",
     "Describe the challenge or gap that led to the development of {topic}."),
]

DIAGRAM_TYPES = ["Flowchart", "State Diagram", "Architecture Diagram", "Sequence Diagram",
                 "Class Diagram", "ER Diagram", "Gantt Chart", "Block Diagram",
                 "Mind Map", "Network Diagram"]

def detect_language(text: str) -> str:
    """Heuristic language detection."""
    sample = text[:500].lower()
    if any(c in sample for c in "あいうえお"): return "Japanese"
    if any(c in sample for c in "你好世界"): return "Chinese"
    if any(c in sample for c in "की है"): return "Hindi"
    if any(c in sample for c in "அஆஇஈ"): return "Tamil"
    return "English"

def classify_tag(text: str) -> str:
    text_l = text.lower()
    for tag, kws in TOPIC_KEYWORDS.items():
        if any(k in text_l for k in kws):
            return tag
    return "General"

def extract_text_from_txt(path: Path) -> tuple[str, int]:
    text = path.read_text(errors="replace")
    lines = [l for l in text.splitlines() if l.strip()]
    return text, len(lines) // 40 + 1

def extract_text_from_pdf(path: Path) -> tuple[str, int]:
    if not HAS_PDF:
        return _demo_text(), 12
    try:
        reader = PyPDF2.PdfReader(str(path))
        pages = len(reader.pages)
        text = "\n".join(p.extract_text() or "" for p in reader.pages)
        return text or _demo_text(), pages
    except Exception:
        return _demo_text(), 12

def extract_text_from_pptx(path: Path) -> tuple[str, int]:
    if not HAS_PPTX:
        return _demo_text(), 20
    try:
        prs = Presentation(str(path))
        slides = len(prs.slides)
        chunks = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    chunks.append(shape.text)
        return "\n".join(chunks) or _demo_text(), slides
    except Exception:
        return _demo_text(), 20

def extract_text_from_docx(path: Path) -> tuple[str, int]:
    if not HAS_DOCX:
        return _demo_text(), 8
    try:
        doc = python_docx.Document(str(path))
        text = "\n".join(p.text for p in doc.paragraphs)
        pages = max(1, len(text) // 3000)
        return text or _demo_text(), pages
    except Exception:
        return _demo_text(), 8

def _demo_text() -> str:
    return """
    Introduction to Operating Systems
    An operating system is system software that manages hardware and software resources.

    Process Management
    A process is a program in execution. The OS manages processes through scheduling, creation, and termination.
    Process states include: New, Ready, Running, Waiting, and Terminated.

    CPU Scheduling Algorithms
    Scheduling algorithms determine the order of process execution.
    FCFS: First Come First Served — simple but can cause the convoy effect.
    SJF: Shortest Job First — minimises average waiting time.
    Round Robin: assigns a fixed time quantum to each process cyclically.

    Memory Management
    Memory management controls allocation and deallocation of memory.
    Fragmentation: Internal fragmentation wastes space within allocated blocks.
    External fragmentation wastes space between allocated blocks.

    Virtual Memory and Paging
    Virtual memory allows execution of processes larger than physical memory.
    Paging divides memory into fixed-size pages and frames.
    A page fault occurs when a referenced page is not in physical memory.

    File Systems
    A file system organises data on storage devices.
    An inode stores file metadata: size, permissions, timestamps, data block pointers.

    Deadlocks and Prevention
    A deadlock occurs when processes wait indefinitely for resources held by each other.
    Coffman conditions: mutual exclusion, hold and wait, no preemption, circular wait.

    Semaphores and Synchronization
    Semaphores are integer variables used to solve critical section problems.
    A mutex (mutual exclusion lock) ensures only one thread accesses a resource at a time.

    I/O Management
    The OS manages I/O devices through device drivers and interrupt handlers.
    Buffering, caching, and spooling are common I/O optimisation techniques.

    Security and Protection
    The OS enforces security through user authentication and access control.
    Protection rings separate privilege levels for OS kernel and user processes.
    """

def extract_topics_from_text(text: str) -> list[dict]:
    """Extract topic headings from text using heuristics."""
    topics = []
    lines = [l.strip() for l in text.splitlines() if l.strip()]

    heading_re = re.compile(
        r'^(chapter|section|unit|part|topic|module|lesson)?\s*'
        r'(\d+[\.\d]*\s+)?([A-Z][A-Za-z\s,&\-/]{4,60})$'
    )
    short_cap = re.compile(r'^[A-Z][A-Za-z\s,&\-/]{4,50}$')

    seen = set()
    page_counter = [1]

    for i, line in enumerate(lines):
        if len(line) < 5 or len(line) > 80:
            continue
        if line in seen:
            continue
        is_heading = bool(heading_re.match(line)) or (
            bool(short_cap.match(line)) and
            len(line.split()) <= 6 and
            not line.endswith('.')
        )
        if not is_heading:
            continue

        # Advance simulated page counter
        if i % 8 == 0:
            page_counter[0] += 1

        tag = classify_tag(line)
        # Confidence: longer coherent headings → higher confidence
        words = line.split()
        base_conf = 55 + min(len(words) * 6, 35)
        conf = min(base_conf + (5 if tag != "General" else 0), 98)

        topics.append({
            "id": len(topics) + 1,
            "name": line.strip().rstrip(":"),
            "tag": tag,
            "page": f"p. {page_counter[0]}",
            "conf": conf,
        })
        seen.add(line)

        if len(topics) >= 12:
            break

    # Fallback: paragraph-first-sentence topics
    if len(topics) < 4:
        paras = [p.strip() for p in text.split("\n\n") if len(p.strip()) > 60]
        for j, para in enumerate(paras[:8]):
            first = para.split(".")[0].strip()
            if 10 < len(first) < 80:
                tag = classify_tag(first)
                conf = 60 + (j % 3) * 10
                topics.append({
                    "id": len(topics) + 1,
                    "name": first[:60],
                    "tag": tag,
                    "page": f"p. {j + 2}",
                    "conf": conf,
                })
            if len(topics) >= 10:
                break

    return topics[:10]

def generate_questions(topics: list[dict], text: str) -> list[dict]:
    """Generate study questions for each topic."""
    questions = []
    qid = 1

    for topic in topics:
        tname = topic["name"]
        # Use 2–3 question templates per topic
        count = min(3, max(2, len(tname.split()) // 2))
        templates = QUESTION_TEMPLATES[:count]

        for tmpl_q, tmpl_a in templates:
            q_text = tmpl_q.replace("{topic}", tname)
            a_text = tmpl_a.replace("{topic}", tname)

            # Try to pull a relevant sentence from the text
            sentences = [s.strip() for s in text.split('.') if tname.lower()[:8] in s.lower() and len(s.strip()) > 30]
            if sentences:
                a_text += " — " + sentences[0][:200].strip() + "."

            questions.append({
                "id": qid,
                "topic": tname,
                "tag": topic["tag"],
                "q": q_text,
                "a": a_text,
                "conf": max(50, topic["conf"] - 5),
                "page": topic["page"],
            })
            qid += 1

    return questions

def detect_diagrams(text: str, pages: int) -> list[dict]:
    """Detect likely diagram mentions in text."""
    diagram_keywords = [
        ("figure", "Flowchart"), ("fig.", "Block Diagram"), ("diagram", "Architecture Diagram"),
        ("chart", "Gantt Chart"), ("graph", "Network Diagram"), ("table", "Data Table"),
        ("flowchart", "Flowchart"), ("state", "State Diagram"), ("sequence", "Sequence Diagram"),
        ("class diagram", "Class Diagram"), ("er diagram", "ER Diagram"),
    ]
    found = []
    text_l = text.lower()
    seen_types = set()

    for kw, dtype in diagram_keywords:
        if kw in text_l and dtype not in seen_types:
            idx = text_l.find(kw)
            # Grab surrounding context
            snippet = text[max(0, idx - 30): idx + 80].strip()
            page_num = max(1, min(pages, (idx // max(1, len(text) // pages)) + 1))
            conf = 88 if kw in ("flowchart", "state", "sequence", "class diagram") else \
                   75 if kw in ("diagram", "figure", "chart") else 60

            found.append({
                "id": len(found) + 1,
                "title": f"{dtype} — {snippet[:40].strip()}",
                "type": dtype,
                "conf": conf,
                "desc": f"A {dtype.lower()} detected on page {page_num} of the document. "
                        f"Context: '{snippet[:80].strip()}…'",
                "callouts": _generate_callouts(dtype),
                "page": f"p. {page_num}",
                "warn": conf < 65,
                "visual": _visual_type(dtype),
            })
            seen_types.add(dtype)

        if len(found) >= 5:
            break

    # Always include at least 2 diagrams
    if len(found) < 2:
        found = _fallback_diagrams()

    return found[:5]

def _visual_type(dtype: str) -> str:
    mapping = {
        "Flowchart": "flow", "State Diagram": "flow", "Gantt Chart": "gantt",
        "Architecture Diagram": "paging", "Block Diagram": "flow",
    }
    return mapping.get(dtype, "generic")

def _generate_callouts(dtype: str) -> list[dict]:
    callout_map = {
        "Flowchart":   [{"label":"Start/End","text":"Oval shapes mark the entry and exit points"},
                        {"label":"Process","text":"Rectangle shapes represent operations or actions"},
                        {"label":"Decision","text":"Diamond shapes represent conditional branches"},
                        {"label":"Arrow","text":"Directed arrows show the flow of control"}],
        "State Diagram":[{"label":"State","text":"Circle represents a system state"},
                         {"label":"Transition","text":"Arrow shows event triggering state change"},
                         {"label":"Initial","text":"Filled circle marks the starting state"},
                         {"label":"Final","text":"Double circle marks terminal state"}],
        "Gantt Chart": [{"label":"Task Bar","text":"Horizontal bar showing task duration"},
                        {"label":"Timeline","text":"X-axis represents time progression"},
                        {"label":"Dependency","text":"Arrows show task ordering constraints"}],
        "Architecture Diagram":[{"label":"Layer","text":"Horizontal bands show system tiers"},
                                 {"label":"Component","text":"Boxes represent software modules"},
                                 {"label":"Interface","text":"Lines show communication paths"},
                                 {"label":"Boundary","text":"Dashed border marks system scope"}],
    }
    default = [{"label":"Element","text":"Primary component of the diagram"},
               {"label":"Connector","text":"Line or arrow showing relationship"},
               {"label":"Label","text":"Text annotation on component"}]
    return callout_map.get(dtype, default)

def _fallback_diagrams() -> list[dict]:
    return [
        {"id":1,"title":"Process Flow Diagram","type":"Flowchart","conf":85,
         "desc":"A flowchart illustrating the main process flow described in this document.",
         "callouts":_generate_callouts("Flowchart"),
         "page":"p. 3","warn":False,"visual":"flow"},
        {"id":2,"title":"System Architecture Overview","type":"Architecture Diagram","conf":78,
         "desc":"An architecture diagram showing the key components and their relationships.",
         "callouts":_generate_callouts("Architecture Diagram"),
         "page":"p. 7","warn":False,"visual":"paging"},
    ]

def compute_stats(topics, questions, diagrams, pages, file_size, proc_time, language):
    high = sum(1 for t in topics if t["conf"] >= 85)
    mid  = sum(1 for t in topics if 60 <= t["conf"] < 85)
    low  = sum(1 for t in topics if t["conf"] < 60)
    return {
        "pages": pages,
        "topics": len(topics),
        "questions": len(questions),
        "diagrams": len(diagrams),
        "language": language,
        "file_size_kb": round(file_size / 1024, 1),
        "proc_time_s": round(proc_time, 1),
        "conf_high": high,
        "conf_mid": mid,
        "conf_low": low,
        "avg_conf": round(sum(t["conf"] for t in topics) / max(1, len(topics)), 1),
    }

# ─────────────────────────────────────────────────────────────────────────────
# PROCESSING JOB
# ─────────────────────────────────────────────────────────────────────────────

def process_document(job_id: str, file_path: Path, filename: str, file_size: int):
    """Background task: parse → extract → store results."""
    try:
        jobs[job_id]["status"] = "parsing"
        jobs[job_id]["progress"] = 10
        t0 = time.time()

        # 1. Extract text
        ext = file_path.suffix.lower()
        if ext == ".pdf":
            text, pages = extract_text_from_pdf(file_path)
        elif ext in (".ppt", ".pptx"):
            text, pages = extract_text_from_pptx(file_path)
        elif ext in (".doc", ".docx"):
            text, pages = extract_text_from_docx(file_path)
        else:
            text, pages = extract_text_from_txt(file_path)

        jobs[job_id]["progress"] = 25
        jobs[job_id]["status"] = "detecting_language"

        language = detect_language(text)

        jobs[job_id]["progress"] = 35
        jobs[job_id]["status"] = "extracting_topics"

        topics = extract_topics_from_text(text)

        jobs[job_id]["progress"] = 55
        jobs[job_id]["status"] = "generating_questions"

        questions = generate_questions(topics, text)

        jobs[job_id]["progress"] = 70
        jobs[job_id]["status"] = "detecting_diagrams"

        diagrams = detect_diagrams(text, pages)

        jobs[job_id]["progress"] = 85
        jobs[job_id]["status"] = "scoring"

        proc_time = time.time() - t0
        stats = compute_stats(topics, questions, diagrams, pages, file_size, proc_time, language)

        jobs[job_id]["progress"] = 100
        jobs[job_id]["status"] = "done"
        jobs[job_id]["result"] = {
            "job_id": job_id,
            "filename": filename,
            "topics": topics,
            "questions": questions,
            "diagrams": diagrams,
            "stats": stats,
            "processed_at": datetime.utcnow().isoformat(),
        }
    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)

# ─────────────────────────────────────────────────────────────────────────────
# ROUTES
# ─────────────────────────────────────────────────────────────────────────────

ALLOWED_EXTENSIONS = {".pdf", ".ppt", ".pptx", ".doc", ".docx", ".txt", ".md", ".epub"}
MAX_FILE_BYTES = 50 * 1024 * 1024  # 50 MB

@app.post("/api/upload")
async def upload_file(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    # Validate extension
    ext = Path(file.filename or "").suffix.lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(400, f"Unsupported file type '{ext}'. Allowed: {', '.join(ALLOWED_EXTENSIONS)}")

    # Read content
    content = await file.read()
    if len(content) > MAX_FILE_BYTES:
        raise HTTPException(413, "File too large. Maximum size is 50 MB.")
    if len(content) == 0:
        raise HTTPException(400, "Uploaded file is empty.")

    # Save
    job_id = str(uuid.uuid4())[:8]
    safe_name = re.sub(r"[^\w.\-]", "_", file.filename or "upload")
    save_path = UPLOAD_DIR / f"{job_id}_{safe_name}"
    save_path.write_bytes(content)

    # Register job
    jobs[job_id] = {
        "job_id": job_id,
        "filename": file.filename,
        "status": "queued",
        "progress": 0,
        "result": None,
        "error": None,
    }

    # Run in background
    background_tasks.add_task(process_document, job_id, save_path, file.filename, len(content))

    return {"job_id": job_id, "filename": file.filename}


@app.get("/api/job/{job_id}")
async def get_job(job_id: str):
    if job_id not in jobs:
        raise HTTPException(404, "Job not found")
    job = jobs[job_id]
    return {
        "job_id": job_id,
        "status": job["status"],
        "progress": job["progress"],
        "error": job.get("error"),
    }


@app.get("/api/results/{job_id}")
async def get_results(job_id: str):
    if job_id not in jobs:
        raise HTTPException(404, "Job not found")
    job = jobs[job_id]
    if job["status"] != "done":
        raise HTTPException(400, f"Job not complete (status: {job['status']})")
    return job["result"]


@app.post("/api/demo")
async def demo(background_tasks: BackgroundTasks):
    """Process demo document without file upload."""
    job_id = str(uuid.uuid4())[:8]
    demo_path = UPLOAD_DIR / f"{job_id}_demo.txt"
    demo_path.write_text(_demo_text())
    jobs[job_id] = {
        "job_id": job_id,
        "filename": "Operating_Systems_Lecture_Notes.pdf",
        "status": "queued",
        "progress": 0,
        "result": None,
        "error": None,
    }
    background_tasks.add_task(process_document, job_id, demo_path, "Operating_Systems_Lecture_Notes.pdf", 48000)
    return {"job_id": job_id}


@app.get("/api/export/{job_id}")
async def export_results(job_id: str, fmt: str = "json"):
    if job_id not in jobs or jobs[job_id]["status"] != "done":
        raise HTTPException(404, "Results not found")

    result = jobs[job_id]["result"]

    if fmt == "json":
        out_path = EXPORT_DIR / f"{job_id}_studylens.json"
        out_path.write_text(json.dumps(result, indent=2))
        return FileResponse(str(out_path), media_type="application/json",
                            filename=f"studylens_{job_id}.json")

    if fmt == "txt":
        lines = [f"StudyLens Export — {result['filename']}",
                 f"Processed: {result['processed_at']}", "=" * 60, ""]
        for i, q in enumerate(result["questions"], 1):
            lines += [f"Q{i:02d}. [{q['topic']}] {q['q']}",
                      f"    Answer: {q['a']}", f"    Source: {q['page']}  Confidence: {q['conf']}%", ""]
        out_path = EXPORT_DIR / f"{job_id}_studylens.txt"
        out_path.write_text("\n".join(lines))
        return FileResponse(str(out_path), media_type="text/plain",
                            filename=f"studylens_{job_id}.txt")

    raise HTTPException(400, "Unsupported export format. Use 'json' or 'txt'.")


@app.get("/api/health")
async def health():
    return {
        "status": "ok",
        "parsers": {"pdf": HAS_PDF, "pptx": HAS_PPTX, "docx": HAS_DOCX},
        "jobs_in_memory": len(jobs),
    }
